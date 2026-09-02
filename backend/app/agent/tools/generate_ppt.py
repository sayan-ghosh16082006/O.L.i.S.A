import time
import asyncio
from pathlib import Path
from typing import List, Optional
from pydantic import BaseModel, Field
from pptx import Presentation
from langchain_core.tools import StructuredTool


class SimplePPTInput(BaseModel):
    deck_title: str = Field(description="The main title of the presentation")
    slides: List[dict] = Field(
        description="List of slides. Each dict needs 'type' ('title' or 'content') and 'title'. Content slides need a 'bullets' list of strings."
    )
    filename: Optional[str] = Field(None, description="Optional name for the file.")


def _generate_pptx_sync(deck_title: str, slides: List[dict], filename: Optional[str] = None) -> str:
    """The actual blocking logic for generating the PPTX."""


    base_dir = Path(__file__).resolve().parent.parent.parent.parent
    workspace_dir = base_dir / "project_workspace"
    workspace_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    for slide_data in slides:
        stype = slide_data.get("type", "content").lower()
        title_text = slide_data.get("title", "Untitled Slide")

        if stype == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_text
            if "subtitle" in slide_data:
                slide.placeholders[1].text = slide_data["subtitle"]
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text
            
            bullets = slide_data.get("bullets", [])
            tf = slide.placeholders[1].text_frame
            tf.clear() 

            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b.get("text", str(b)) if isinstance(b, dict) else str(b)

    if not filename:
        clean_title = "".join(c for c in deck_title if c.isalnum() or c in " _-").strip()
        filename = f"{clean_title or 'deck'}_{int(time.time())}.pptx"
    
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    save_path = workspace_dir / filename
    prs.save(str(save_path))
    
    return f"Successfully saved to: {save_path.absolute()}"


async def agenerate_pptx_simple(deck_title: str, slides: List[dict], filename: Optional[str] = None) -> str:
    """Async wrapper that runs the sync PPTX generation in a separate thread."""
    return await asyncio.to_thread(_generate_pptx_sync, deck_title, slides, filename)

# LangChain Tool Integration (Async)
create_pptx_tool = StructuredTool.from_function(
    func=_generate_pptx_sync,         
    coroutine=agenerate_pptx_simple, 
    name="generate_pptx",
    description=(
        "Generates a PowerPoint file from content and saves it to the workspace. "
        "Input 'slides' should be a list of dicts. "
        "Example: [{'type': 'title', 'title': 'Intro', 'subtitle': 'by AI'}, "
        "{'type': 'content', 'title': 'Points', 'bullets': ['point 1', 'point 2']}]"
    ),
    args_schema=SimplePPTInput
)